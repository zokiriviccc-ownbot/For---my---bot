from flask import Flask
from threading import Thread
import os

app = Flask('')

@app.route('/')
def home():
    return "Bot ishlayapti!"

def run():
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 8080)))

def keep_alive():
    t = Thread(target=run)
    t.start()
import telebot
from telebot import types
import google.generativeai as genai
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import requests
from io import BytesIO

# -----------------------------
# TOKEN & API KEY
# -----------------------------
TOKEN = "7686001116:AAGL7ynDrm3TedLAyDWJaDln2Jv_-LNRlhM"
GEMINI_KEY = "AIzaSyAypYbPDfJzQAL3YmJcTl6NyUGwV6P61e4"

CARD_NUMBER = "9860 0101 0421 3771"  # Premium karta
PRICE = "3000 so'm"
FREE_LIMIT = 3
# -----------------------------

bot = telebot.TeleBot(TOKEN)
genai.configure(api_key=GEMINI_KEY)
model = genai.GenerativeModel("gemini-pro")

user_limits = {}
user_data = {}

# /start
@bot.message_handler(commands=['start'])
def start(message):
    uid = message.from_user.id
    user_limits.setdefault(uid,0)

    markup = types.InlineKeyboardMarkup()
    btn1 = types.InlineKeyboardButton("📄 Referat va PPT yaratish",callback_data="create")
    btn2 = types.InlineKeyboardButton("💎 Premium olish",callback_data="premium")
    markup.add(btn1,btn2)

    left = FREE_LIMIT - user_limits[uid]
    text = f"👋 Salom!\nSizda {left} ta bepul imkoniyat bor"
    bot.send_message(message.chat.id,text,reply_markup=markup)

# Tugma bosish
@bot.callback_query_handler(func=lambda call: True)
def callback(call):
    uid = call.from_user.id

    # Referat yaratish
    if call.data=="create":
        if user_limits.get(uid,0)>=FREE_LIMIT:
            text=f"""
❌ Sizning bepul limit tugadi

💎 Premium narxi: {PRICE}

💳 To'lov uchun karta:
{CARD_NUMBER}

To'lov qilgandan keyin chek yuboring.
"""
            bot.send_message(call.message.chat.id,text)
            return
        msg = bot.send_message(call.message.chat.id,"📚 Referat mavzusini yozing:")
        bot.register_next_step_handler(msg,get_topic)

    # Premium
    if call.data=="premium":
        text=f"""
💎 Premium

Narx: {PRICE}

💳 To'lov uchun karta:
{CARD_NUMBER}

To'lov qilgandan keyin chek yuboring.
"""
        bot.send_message(call.message.chat.id,text)

# Mavzu
def get_topic(message):
    user_data[message.chat.id] = {}
    user_data[message.chat.id]["topic"] = message.text
    msg = bot.send_message(message.chat.id,"🏫 Universitet nomini yozing:")
    bot.register_next_step_handler(msg,get_uni)

# Universitet
def get_uni(message):
    user_data[message.chat.id]["uni"] = message.text
    msg = bot.send_message(message.chat.id,"👤 Ism familiya:")
    bot.register_next_step_handler(msg,get_name)

# Talaba
def get_name(message):
    user_data[message.chat.id]["name"] = message.text
    msg = bot.send_message(message.chat.id,"👨‍🏫 O‘qituvchi ismi:")
    bot.register_next_step_handler(msg,get_teacher)

# O'qituvchi
def get_teacher(message):
    user_data[message.chat.id]["teacher"] = message.text
    bot.send_message(message.chat.id,"⏳ Referat va PPT tayyorlanmoqda...")
    generate(message)

# AI generatsiya
def generate(message):
    uid = message.from_user.id
    data = user_data[uid]
    topic = data["topic"]

    # AI matn yaratish
    response = model.generate_content(f"{topic} haqida talaba uchun referat yoz")
    text = response.text

    # Word va PPT yaratish
    create_word(message.chat.id,topic,text)
    create_ppt_with_images(message.chat.id,topic,text)

    # Limit oshirish
    user_limits[uid] += 1

# Word fayl yaratish
def create_word(chat_id,topic,text):
    doc = Document()
    doc.add_heading(topic,0)
    doc.add_paragraph(text)
    file="referat.docx"
    doc.save(file)
    bot.send_document(chat_id,open(file,"rb"))

# PPT + rasm
def create_ppt_with_images(chat_id,topic,text):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    # 1-slayd
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = topic
    slide.placeholders[1].text = text[:500]

    try:
        url = f"https://source.unsplash.com/600x400/?{topic.replace(' ','')}"
        response = requests.get(url)
        image_stream = BytesIO(response.content)
        slide.shapes.add_picture(image_stream,Inches(1),Inches(2.5),width=Inches(6))
    except:
        pass

    # Keyingi slaydlar (textni bo‘lish)
    parts = [text[i:i+500] for i in range(500,len(text),500)]
    for part in parts:
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = topic
        slide.placeholders[1].text = part
        try:
            url = f"https://source.unsplash.com/600x400/?{topic.replace(' ','')}"
            response = requests.get(url)
            image_stream = BytesIO(response.content)
            slide.shapes.add_picture(image_stream,Inches(1),Inches(2.5),width=Inches(6))
        except:
            pass

    file="slayd_rasmli.pptx"
    prs.save(file)
    bot.send_document(chat_id,open(file,"rb"))
keep_alive()

bot.polling()
